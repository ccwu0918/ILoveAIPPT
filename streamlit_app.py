import os
import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
import datetime
import google.generativeai as genai
from dotenv import load_dotenv

# 加載環境變量並配置Google Generative AI
load_dotenv()
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))
model = genai.GenerativeModel('gemini-pro')

def convert_markdown_to_pptx(markdown_text, slide_format, font_size, font_family):
    if not markdown_text.strip():
        st.warning("No content to convert!")
        return None

    ppt = Presentation()

    if slide_format == "Standard (4:3)":
        ppt.slide_width = Inches(10)
        ppt.slide_height = Inches(7.5)
    else:  # Default to Widescreen (16:9)
        ppt.slide_width = Inches(13.33)
        ppt.slide_height = Inches(7.5)

    font_size = Pt(int(font_size))
    lines = markdown_text.splitlines()
    slide = None
    text_frame = None
    current_section = []
    current_section_title = ""

    def add_section_to_slide(section_lines, title=""):
        nonlocal slide, text_frame
        if not section_lines:
            return

        slide = ppt.slides.add_slide(ppt.slide_layouts[5])
        slide.shapes.title.text = title
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
        text_frame = textbox.text_frame
        
        for line in section_lines:
            if line.strip():
                p = text_frame.add_paragraph()
                p.text = line.strip()
                p.font.size = font_size
                p.font.name = font_family

    for line in lines:
        if line.startswith("##"):
            if current_section:
                add_section_to_slide(current_section, current_section_title)
                current_section = []
            current_section_title = line.lstrip("#").strip()
        elif line.startswith("#"):
            current_section_title = line.lstrip("#").strip()
        else:
            current_section.append(line)

    if current_section:
        add_section_to_slide(current_section, current_section_title)

    now = datetime.datetime.now()
    formatted_datetime = now.strftime("%Y-%m-%d-%H-%M-%S")
    save_path = f"{formatted_datetime}-Presentation.pptx"
    ppt.save(save_path)
    
    return save_path

def main():
    st.set_page_config(page_title="I LOVE PPT & SQL Generator", page_icon=":memo:")
    
    st.markdown(
        """
        <div style="text-align:center;">
            <h1>I LOVE PPT & SQL Generator</h1>
            <p>Created by ECF、CTF and CCWu</p>
        </div>
        """,
        unsafe_allow_html=True
    )

    # PPTX 生成部分
    st.header("Markdown to PPTX Converter")
    
    markdown_text = st.text_area("Markdown Input", height=300, placeholder="Enter Markdown here...")
    slide_format = st.selectbox("Slide Format", ["Widescreen (16:9)", "Standard (4:3)"], index=0)
    font_size = st.selectbox("Font Size", [str(i) for i in range(8, 97)], index=10)
    font_family = st.selectbox("Font Family", [
        "Calibri", "Arial", "Times New Roman", "Verdana", "Courier New", "Georgia",
        "Garamond", "Comic Sans MS", "Trebuchet MS", "Palatino Linotype",
        "Tahoma", "Lucida Sans Unicode", "Microsoft Sans Serif",
        "標楷體", "新細明體", "宋体", "黑体", "SimSun", "SimHei",
        "Meiryo", "Malgun Gothic"
    ], index=0)

    if st.button("Convert to PPTX"):
        save_path = convert_markdown_to_pptx(markdown_text, slide_format, font_size, font_family)
        
        if save_path:
            with open(save_path, "rb") as file:
                btn = st.download_button(
                    label="Download Presentation",
                    data=file,
                    file_name=save_path,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )

    # SQL 查詢生成部分
    st.header("SQL Query Generator")
    
    input_query = st.text_area("Enter your query here in Plain Language:")

    if st.button("Generate SQL Query"):
        template = f"""
        Create a SQL query snippet using the below text:
        ```
        {input_query}
        ```
        I just want a SQL Query.
        """
        
        response = model.generate_content(template)
        
        sql_query = response.text.strip().lstrip("```sql").rstrip("```")

        with st.container():
            st.success("SQL query generated successfully:")
            st.code(sql_query, language="sql")

            expected_output_template = f"""
            What would be the expected response of this SQL query snippet:
            ```sql
            {sql_query}
            ```
            Provide sample table Response:
            """
            
            expected_response = model.generate_content(expected_output_template)
            
            st.success("Expected Output of this SQL query will be:")
            st.markdown(expected_response.text)

if __name__ == '__main__':
    main()
